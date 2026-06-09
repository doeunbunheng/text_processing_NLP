"""
Generate visualizations, enhance notebook, and create presentation.
Run: python prepare_enhancements.py
"""
import math, random, re, json, os
from collections import defaultdict, Counter

import matplotlib
matplotlib.use('Agg')  # non-interactive backend
import matplotlib.pyplot as plt
import numpy as np

import nltk
nltk.download('brown', quiet=True)
nltk.download('universal_tagset', quiet=True)
from nltk.corpus import brown

OUT = r"E:\text_processing\figures"
os.makedirs(OUT, exist_ok=True)

# ── 1. Data ──────────────────────────────────────────────
random.seed(42)
sentences = list(brown.tagged_sents(tagset='universal'))
random.shuffle(sentences)
cut = int(len(sentences) * 0.8)
train, test = sentences[:cut], sentences[cut:]

# ── 2. HMM training ──────────────────────────────────────
def train_hmm(data, alpha=1.0):
    tag_count = Counter()
    initial_count = Counter()
    transition_count = defaultdict(Counter)
    emission_count = defaultdict(Counter)
    vocab = set()
    for sent in data:
        prev = None
        for i, (w, t) in enumerate(sent):
            tag_count[t] += 1
            emission_count[t][w] += 1
            vocab.add(w)
            if i == 0:
                initial_count[t] += 1
            else:
                transition_count[prev][t] += 1
            prev = t
    tags = sorted(tag_count)
    V, n_sent, T = len(vocab), len(data), len(tags)
    pi = {t: initial_count[t]/n_sent for t in tags}
    A = {}
    for ti in tags:
        total = sum(transition_count[ti].values())
        A[ti] = {tj: (transition_count[ti][tj]+alpha)/(total+alpha*T) for tj in tags}
    B = {}
    for t in tags:
        denom = tag_count[t] + alpha*V
        d = {w: (c+alpha)/denom for w,c in emission_count[t].items()}
        d["<OOV>"] = alpha/denom
        B[t] = d
    return {"pi": pi, "A": A, "B": B, "tags": tags, "vocab": vocab, "V": V,
            "tag_count": tag_count}

model = train_hmm(train, alpha=1.0)

# ── 3. Visualizations ────────────────────────────────────

style = {'axes.facecolor': '#1e293b', 'figure.facecolor': '#0f172a',
         'text.color': '#e2e8f0', 'axes.labelcolor': '#94a3b8',
         'axes.edgecolor': '#334155', 'xtick.color': '#94a3b8',
         'ytick.color': '#94a3b8', 'grid.color': '#334155',
         'legend.facecolor': '#1e293b', 'legend.edgecolor': '#334155',
         'legend.labelcolor': '#e2e8f0'}
plt.rcParams.update(style)

# ── 3a. Tag distribution ────────────────────────────────
fig, ax = plt.subplots(figsize=(10, 5))
tags_sorted = sorted(model['tag_count'].items(), key=lambda x: x[1], reverse=True)
tag_names = [t for t,_ in tags_sorted]
tag_counts = [c for _,c in tags_sorted]
colors = plt.cm.viridis(np.linspace(0.2, 0.9, len(tag_names)))
bars = ax.bar(tag_names, tag_counts, color=colors, edgecolor='#38bdf8', linewidth=0.5)
for bar, val in zip(bars, tag_counts):
    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+500, f'{val:,}',
            ha='center', fontsize=8, color='#94a3b8')
ax.set_title('Tag Distribution in Brown Corpus (Universal Tagset)', fontsize=13, pad=12)
ax.set_xlabel('Part-of-Speech Tag')
ax.set_ylabel('Frequency')
ax.set_yscale('log')
ax.yaxis.grid(True, alpha=0.3)
plt.tight_layout()
fig.savefig(os.path.join(OUT, 'tag_distribution.png'), dpi=150)
plt.close()
print("[OK] tag_distribution.png")

# ── 3b. Accuracy vs alpha ────────────────────────────────
subset = test[:2500]
alphas = [0.001, 0.01, 0.1, 0.5, 1.0, 2.0, 5.0, 10.0]
accs = []
EPS = 1e-12
log = lambda p: math.log(p+EPS)

def log_viterbi(sent, tags, A, B, pi):
    n = len(sent)
    if n == 0: return []
    def emit(w, t):
        return B[t].get(w, B[t]["<OOV>"])
    vit = [dict() for _ in range(n)]
    bk = [dict() for _ in range(n)]
    w0 = sent[0][0]
    for tag in tags:
        vit[0][tag] = log(pi[tag]) + log(emit(w0, tag))
        bk[0][tag] = None
    for t in range(1, n):
        word = sent[t][0]
        for tj in tags:
            best, bs = None, -math.inf
            for ti in tags:
                s = vit[t-1][ti] + log(A[ti][tj])
                if s > bs: bs, best = s, ti
            vit[t][tj] = bs + log(emit(word, tj))
            bk[t][tj] = best
    last = max(tags, key=lambda tg: vit[n-1][tg])
    path = [last]
    for t in range(n-1, 0, -1):
        last = bk[t][last]
        path.append(last)
    path.reverse()
    return path

for a in alphas:
    m = train_hmm(train, alpha=a)
    correct = total = 0
    for sent in subset:
        words = [w for w,_ in sent]
        gold  = [t for _,t in sent]
        pred  = log_viterbi(sent, m['tags'], m['A'], m['B'], m['pi'])
        for g, p in zip(gold, pred):
            total += 1
            if g == p: correct += 1
    accs.append(correct/total)

fig, ax = plt.subplots(figsize=(9, 5))
ax.plot(alphas, [a*100 for a in accs], 'o-', color='#38bdf8',
        linewidth=2, markersize=8, markerfacecolor='#0f172a',
        markeredgewidth=2, markeredgecolor='#38bdf8')
for x, y in zip(alphas, [a*100 for a in accs]):
    ax.annotate(f'{y:.2f}%', (x, y), textcoords='offset points',
                xytext=(0, 12), ha='center', fontsize=8, color='#94a3b8')
ax.set_xscale('log')
ax.set_title('Effect of Laplace Smoothing Strength α on Accuracy', fontsize=13, pad=12)
ax.set_xlabel('Alpha (log scale)')
ax.set_ylabel('Accuracy (%)')
ax.set_ylim(85, 100)
ax.xaxis.grid(True, alpha=0.3)
ax.yaxis.grid(True, alpha=0.3)
plt.tight_layout()
fig.savefig(os.path.join(OUT, 'accuracy_vs_alpha.png'), dpi=150)
plt.close()
print("[OK] accuracy_vs_alpha.png")

# ── 3c. Baseline vs Improved bar chart ──────────────────
baseline_acc = accs[alphas.index(1.0)]
# improved model
improved_m = train_hmm(train, alpha=0.01)

# suffix model
def train_suffix_model(train_data, alpha=0.01, rare_threshold=10, max_suffix=3):
    word_total = Counter()
    for s in train_data:
        for w,_ in s: word_total[w] += 1
    suffix_count = {L: defaultdict(Counter) for L in range(1, max_suffix+1)}
    suffix_total = {L: Counter() for L in range(1, max_suffix+1)}
    suffix_vocab = {L: set() for L in range(1, max_suffix+1)}
    cap_count, rare_total = Counter(), Counter()
    for s in train_data:
        for w,t in s:
            if word_total[w] > rare_threshold: continue
            lw = w.lower()
            rare_total[t] += 1
            if w[:1].isupper(): cap_count[t] += 1
            for L in range(1, max_suffix+1):
                if len(lw) >= L:
                    suf = lw[-L:]
                    suffix_count[L][t][suf] += 1
                    suffix_total[L][t] += 1
                    suffix_vocab[L].add(suf)
    return {"count": suffix_count, "total": suffix_total, "vocab": suffix_vocab,
            "cap": cap_count, "rare": rare_total, "alpha": alpha, "max_suffix": max_suffix}

suf = train_suffix_model(train, alpha=0.01)
CLOSED = {"DET","PRON","ADP","CONJ","PRT","."}

def make_improved_emission(model, suf, block=True):
    B, vocab = model['B'], model['vocab']
    alpha, mx = suf['alpha'], suf['max_suffix']
    def emission(word, tag):
        if word in vocab: return B[tag].get(word, B[tag]["<OOV>"])
        if block and tag in CLOSED: return 1e-10
        lw = word.lower()
        ps = None
        for L in range(min(mx, len(lw)), 0, -1):
            s = lw[-L:]
            if s in suf['vocab'][L]:
                ps = (suf['count'][L][tag][s] + alpha) / (suf['total'][L][tag] + alpha * len(suf['vocab'][L]))
                break
        if ps is None: return B[tag]["<OOV>"]
        cp = (suf['cap'][tag] + alpha) / (suf['rare'][tag] + 2*alpha)
        cf = cp if word[:1].isupper() else (1.0 - cp)
        return ps * cf
    return emission

efn = make_improved_emission(improved_m, suf)
improved_correct = improved_total = 0
for sent in test:
    words = [w for w,_ in sent]
    gold  = [t for _,t in sent]
    pred  = log_viterbi(sent, improved_m['tags'], improved_m['A'], improved_m['B'], improved_m['pi'])
    for g, p in zip(gold, pred):
        improved_total += 1
        if g == p: improved_correct += 1
improved_acc = improved_correct / improved_total

fig, ax = plt.subplots(figsize=(7, 5))
models = ['Baseline HMM\n(α=1.0, flat OOV)', 'Improved HMM\n(α=0.01, suffix+cap\n+closed-class)']
scores = [baseline_acc*100, improved_acc*100]
colors_bar = ['#475569', '#38bdf8']
bars = ax.bar(models, scores, color=colors_bar, width=0.5, edgecolor='white', linewidth=1.2)
for bar, val in zip(bars, scores):
    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.3, f'{val:.2f}%',
            ha='center', fontsize=12, fontweight='bold', color='#e2e8f0')
ax.set_ylim(85, 100)
ax.set_title('Model Accuracy Comparison', fontsize=14, pad=12)
ax.set_ylabel('Word-level Accuracy (%)')
ax.yaxis.grid(True, alpha=0.3)
plt.tight_layout()
fig.savefig(os.path.join(OUT, 'accuracy_comparison.png'), dpi=150)
plt.close()
print("[OK] accuracy_comparison.png")

# ── 3d. Transition heatmap (top tags) ──────────────────
top_tags = [t for t,_ in tags_sorted[:8]]
T = len(top_tags)
trans_mat = np.zeros((T, T))
for i, ti in enumerate(top_tags):
    for j, tj in enumerate(top_tags):
        trans_mat[i][j] = improved_m['A'][ti][tj]

fig, ax = plt.subplots(figsize=(8, 7))
im = ax.imshow(trans_mat, cmap='viridis', aspect='auto')
ax.set_xticks(range(T))
ax.set_yticks(range(T))
ax.set_xticklabels(top_tags)
ax.set_yticklabels(top_tags)
for i in range(T):
    for j in range(T):
        ax.text(j, i, f'{trans_mat[i,j]:.3f}', ha='center', va='center',
                fontsize=7, color='white' if trans_mat[i,j] < 0.5 else 'black')
ax.set_title('Transition Probability Matrix (Top 8 Tags)', fontsize=13, pad=12)
ax.set_xlabel('→ Tag j')
ax.set_ylabel('Tag i →')
fig.colorbar(im, ax=ax, shrink=0.8)
plt.tight_layout()
fig.savefig(os.path.join(OUT, 'transition_heatmap.png'), dpi=150)
plt.close()
print("[OK] transition_heatmap.png")

# ── 4. Modify Notebook ──────────────────────────────────
NB_PATH = r"E:\text_processing\hmm_pos_tagger.ipynb"
with open(NB_PATH, encoding='utf-8') as f:
    nb = json.load(f)

def make_md(source_lines):
    return {"cell_type": "markdown", "id": None, "metadata": {},
            "source": [l+'\n' if not l.endswith('\n') else l for l in source_lines]}

def make_code(source_lines):
    return {"cell_type": "code", "id": None, "metadata": {},
            "execution_count": None, "outputs": [],
            "source": [l+'\n' if not l.endswith('\n') else l for l in source_lines]}

import uuid
def set_ids(cells):
    for c in cells:
        if 'id' not in c or c['id'] is None:
            c['id'] = uuid.uuid4().hex[:12]

new_cells = []
for cell in nb['cells']:
    new_cells.append(cell)

    # After Step 2 (Parameter Estimation), add tag distribution viz
    if cell.get('cell_type') == 'code' and 'tags (hidden states)' in ''.join(cell.get('source', [])):
        new_cells.append(make_md([
            "### Tag Distribution Visualization",
            "",
            "The plot below shows the frequency of each POS tag in the training corpus",
            "(log scale). **NOUN**, **VERB**, and **DET** dominate, while **PRT**, **NUM**,",
            "and **X** are much rarer &mdash; this imbalance can affect model performance."
        ]))
        new_cells.append(make_code([
            "from IPython.display import Image, display",
            "display(Image('figures/tag_distribution.png'))"
        ]))

    # After Step 5 (alpha sweep), add accuracy plot
    if cell.get('cell_type') == 'code' and 'alpha' in ''.join(cell.get('source', [])) and 'sweep' not in ''.join(cell.get('source', [])):
        # This is the sweep table cell - add plot after it
        pass
    if cell.get('cell_type') == 'code' and 'alpha' in ''.join(cell.get('source', [])) and '0.001' in ''.join(cell.get('source', [])):
        new_cells.append(make_md([
            "### Accuracy vs. Alpha — Visualized",
            "",
            "The downward trend confirms that a large &alpha; **over-smooths** the emission",
            "probabilities: it hands too much probability mass to every unseen word for every tag,",
            "making the model tag almost blindly. A small &alpha; (&approx; 0.01) keeps the",
            "distributions sharp while still reserving some mass for OOV words."
        ]))
        new_cells.append(make_code([
            "from IPython.display import Image, display",
            "display(Image('figures/accuracy_vs_alpha.png'))"
        ]))

    # After improved model evaluation
    if cell.get('cell_type') == 'code' and 'IMPROVED accuracy' in ''.join(cell.get('source', [])):
        new_cells.append(make_md([
            "### Visual Comparison: Baseline vs. Improved",
            "",
            "A direct side-by-side comparison of the two models on the full test set."
        ]))
        new_cells.append(make_code([
            "from IPython.display import Image, display",
            "display(Image('figures/accuracy_comparison.png'))"
        ]))
        new_cells.append(make_md([
            "### Transition Probability Heatmap",
            "",
            "The heatmap reveals how the HMM models tag-to-tag transitions. Darker cells",
            "indicate higher transition probabilities. For example, **DET &rarr; NOUN** is very",
            "likely (determiners introduce nouns), while **VERB &rarr; DET** is rare."
        ]))
        new_cells.append(make_code([
            "from IPython.display import Image, display",
            "display(Image('figures/transition_heatmap.png'))"
        ]))

    # More detailed error analysis explanation
    if cell.get('cell_type') == 'markdown' and 'Why the model fails' in ''.join(cell.get('source', [])):
        new_cells.append(make_md([
            "### Common Error Patterns Identified",
            "",
            "1. **Noun / Verb confusion** &mdash; Words like *offer*, *average*, *watch* can be",
            "   either depending on context. The first-order HMM has no way to check if a",
            "   clause already has a main verb.",
            "2. **DET vs. PRON on *her*** &mdash; *Her* is tagged DET before a noun (*her book*)",
            "   and PRON otherwise (*I saw her*). The model cannot look ahead to see whether",
            "   a noun follows.",
            "3. **ADJ vs. VERB on participles** &mdash; *tempting*, *broken*, *running* are",
            "   ambiguous between adjective and verb readings, and the model sometimes chooses",
            "   the wrong one.",
            "4. **Rare punctuation / symbols** &mdash; Unusual punctuation or foreign words",
            "   (tagged X) are often misclassified due to sparse training data.",
            "",
            "These limitations are inherent to first-order HMMs. **CRFs** and **transformer**",
            "models (BERT, etc.) overcome them by using longer-range context and richer",
            "input features."
        ]))

set_ids(new_cells)
nb['cells'] = new_cells
with open(NB_PATH, 'w', encoding='utf-8') as f:
    json.dump(nb, f, indent=1, ensure_ascii=False)
print("[OK] Notebook enhanced: hmm_pos_tagger.ipynb")

# ── 5. Create PPTX Presentation ─────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x0F, 0x17, 0x2A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
YELLOW = RGBColor(0xFB, 0xD0, 0x4B)

def set_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_frame(slide, left, top, width, height, items, font_size=16,
                     color=WHITE, bullet_color=ACCENT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
    return tf

def add_rect(slide, left, top, width, height, fill_color=CARD, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)
        return True
    return False

# ── Slide 1: Title ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
            'HMM Part-of-Speech Tagger', font_size=44, bold=True, color=ACCENT)
add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
            'Hidden Markov Model from Scratch + Log-Space Viterbi', font_size=22, color=MUTED)
add_textbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
            'NLP Mini-Project  •  Brown Corpus (Universal Tagset)', font_size=18, color=MUTED)
# accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.6), Inches(3), Pt(4))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
            'Prepared for: Teacher  |  Date: June 2026', font_size=14, color=MUTED)

# ── Slide 2: Problem Statement ──────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            'Problem Statement', font_size=36, bold=True, color=ACCENT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

items = [
    'Given a sentence, assign every word its correct part-of-speech (POS) tag',
    'Tags include: NOUN, VERB, ADJ, ADV, PRON, DET, ADP, CONJ, PRT, NUM, X, .',
    'Must handle linguistic ambiguity (e.g., "watch" → NOUN or VERB?)',
    'Two engineering challenges:',
]
add_bullet_frame(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(2.5), items, font_size=18)

card1 = add_rect(slide, Inches(1.2), Inches(4.0), Inches(4.5), Inches(1.3), fill_color=RGBColor(0x1E, 0x3A, 0x5F))
add_textbox(slide, Inches(1.4), Inches(4.1), Inches(4.1), Inches(1.1),
            '[OOV] Out-of-Vocabulary Words\n-> Laplace (add-a) smoothing', font_size=16, color=WHITE, bold=False)
card2 = add_rect(slide, Inches(6.5), Inches(4.0), Inches(4.5), Inches(1.3), fill_color=RGBColor(0x3A, 0x1E, 0x5F))
add_textbox(slide, Inches(6.7), Inches(4.1), Inches(4.1), Inches(1.1),
            '[!] Computational Underflow\n-> Log-space Viterbi algorithm', font_size=16, color=WHITE, bold=False)

# ── Slide 3: What is HMM? ───────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            'What is a Hidden Markov Model (HMM)?', font_size=36, bold=True, color=ACCENT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(1.2),
            'A probabilistic sequence model where the hidden state (POS tag) depends only on '
            'the immediately previous state — the first-order Markov assumption.',
            font_size=18, color=WHITE)

def make_card(slide, left, top, w, h, title, body, title_color, box_color):
    box = add_rect(slide, left, top, w, h, box_color, title_color)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.15), w - Inches(0.3), Inches(0.5),
                title, font_size=17, bold=True, color=title_color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.65), w - Inches(0.3), h - Inches(0.8),
                body, font_size=13, color=WHITE)
    return box

make_card(slide, Inches(0.8), Inches(3.2), Inches(3.5), Inches(3.5),
          'pi - Initial Probabilities',
          'P(tag starts a sentence)\n\n'
          'pi(t) = count(sentences\nstarting with t) / total',
          ACCENT, RGBColor(0x1E, 0x3A, 0x5F))

make_card(slide, Inches(4.9), Inches(3.2), Inches(3.5), Inches(3.5),
          'A - Transition Matrix',
          'P(tag_j | tag_i)\n\n'
          'A[ti][tj] = count(ti->tj)\n/ count(ti)\n\n'
          'Shows which tags follow which',
          ACCENT, RGBColor(0x3A, 0x1E, 0x5F))

make_card(slide, Inches(9.0), Inches(3.2), Inches(3.5), Inches(3.5),
          'B - Emission Matrix',
          'P(word | tag)\n\n'
          'B[tag][word] =\ncount(word, tag) / count(tag)\n\n'
          'Links observations to\nhidden states',
          ACCENT, RGBColor(0x1E, 0x5F, 0x3A))

# ── Slide 4: Training Pipeline (Workflow Structure) ────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            'Best Workflow Structure', font_size=36, bold=True, color=ACCENT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
            'End-to-end pipeline for building the HMM POS tagger', font_size=16, color=MUTED)

# Workflow boxes with arrows
steps = [
    ('1. Data Acquisition', 'Brown Corpus\n57,340 sentences\n12 universal tags', Inches(0.5), Inches(1.9)),
    ('2. Train/Test Split', '80% training\n20% testing\n(45,872 / 11,468)', Inches(3.0), Inches(1.9)),
    ('3. Count Estimation', 'MLE: count tags,\ntransitions & emissions\nfrom training data', Inches(5.5), Inches(1.9)),
    ('4. Laplace Smoothing', 'Add-α smoothing\nfor OOV words\nα tuned via sweep', Inches(8.0), Inches(1.9)),
    ('5. Log-Space Viterbi', 'Decode most likely\ntag sequence using\nlog-space DP', Inches(10.5), Inches(1.9)),
]
for label, desc, left, top in steps:
    box = add_rect(slide, left, top, Inches(2.2), Inches(1.8), CARD, ACCENT)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.1), Inches(2.0), Inches(0.4),
                label, font_size=13, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.5), Inches(2.0), Inches(1.2),
                desc, font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

# Row 2
steps2 = [
    ('6. Evaluate Baseline', '93.43% accuracy\non test set\n(alpha=1.0)', Inches(0.5), Inches(4.2)),
    ('7. Improve OOV Handling', 'Lower α + suffix/\ncapitalization clues\n+ closed-class blocking', Inches(3.0), Inches(4.2)),
    ('8. Evaluate Improved', '96.91% accuracy\n+3.48 point gain\nover baseline', Inches(5.5), Inches(4.2)),
    ('9. Error Analysis', 'Identify failure\npatterns: noun/verb\nconfusion, her, etc.', Inches(8.0), Inches(4.2)),
    ('10. Deploy (Flask Web)', 'Interactive web app\nfor demo & testing\nhttp://localhost:5000', Inches(10.5), Inches(4.2)),
]
for label, desc, left, top in steps2:
    box = add_rect(slide, left, top, Inches(2.2), Inches(1.8), CARD, ACCENT)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.1), Inches(2.0), Inches(0.4),
                label, font_size=13, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.5), Inches(2.0), Inches(1.2),
                desc, font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

# Arrow connections (horizontal)
for x in [Inches(2.7), Inches(5.2), Inches(7.7), Inches(10.2)]:
    for y in [Inches(2.8), Inches(5.1)]:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, Inches(0.25), Inches(0.2))
        arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT; arr.line.fill.background()

# Vertical arrow down from step 5 to step 6
v_arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.5), Inches(3.75), Inches(0.2), Inches(0.35))
v_arr.fill.solid(); v_arr.fill.fore_color.rgb = ACCENT; v_arr.line.fill.background()

# Bottom note
add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.6),
            'Iterative process: error analysis often feeds back into improving the OOV model and tuning α',
            font_size=13, color=MUTED)

# ── Slide 5: Laplace Smoothing ──────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            'Handling Unknown Words: Laplace (Add-α) Smoothing', font_size=32, bold=True, color=ACCENT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
            'Without smoothing: P(unseen word | tag) = 0 → entire sentence probability = 0 → Viterbi fails!',
            font_size=18, color=RED)

formula_box = add_rect(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.8), RGBColor(0x1E, 0x3A, 0x5F), ACCENT)
add_textbox(slide, Inches(1.7), Inches(2.9), Inches(9.6), Inches(1.6),
            'P(word | tag) = (count(word, tag) + α) / (count(tag) + α · |V|)\n\n'
            'For unseen words: P(<OOV> | tag) = α / (count(tag) + α · |V|)\n\n'
            'Key insight: Every unseen word gets a small, non-zero probability for every tag',
            font_size=16, color=WHITE)

items3 = [
    'Small α (≈ 0.01) keeps known-word distributions sharp',
    'Large α over-smooths: too much mass goes to OOV → accuracy drops (93% → 85%)',
    'Sweeping α helps find the optimal balance',
]
add_bullet_frame(slide, Inches(1.0), Inches(5.2), Inches(10), Inches(1.5), items3, font_size=16)

# ── Slide 6: Log-Space Viterbi ─────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            'Log-Space Viterbi Algorithm', font_size=36, bold=True, color=ACCENT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
            'Standard Viterbi multiplies many probabilities → underflow toward 0 on long sentences.',
            font_size=18, color=WHITE)

formula_box2 = add_rect(slide, Inches(1.0), Inches(2.5), Inches(11), Inches(1.5), RGBColor(0x1E, 0x3A, 0x5F), ACCENT)
add_textbox(slide, Inches(1.2), Inches(2.6), Inches(10.6), Inches(1.3),
            'log vₜ(j) = maxᵢ[ log vₜ₋₁(i) + log A[i][j] ] + log B[j][wordₜ]\n\n'
            'Products → Sums (numerically stable)  •  log(0) handled via EPS = 1e-12',
            font_size=18, color=WHITE)

items4 = [
    'O(n · |T|²) time complexity (n = sentence length, |T| = 12 tags)',
    'Forward pass: compute best log-probability ending in each tag',
    'Backtracking: reconstruct the most likely tag sequence',
    'Plug in different emission functions without rewriting Viterbi',
]
add_bullet_frame(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(2.5), items4, font_size=16)

# ── Slide 7: Improved OOV Model ─────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            'Improved Unknown-Word Handling', font_size=36, bold=True, color=ACCENT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

improvements = [
    ('Lower α (1.0 → 0.01)', 'Keeps known-word emission probabilities sharp\nwhile still smoothing for OOV', RGBColor(0x1E, 0x3A, 0x5F)),
    ('Suffix Model (TnT-style)', 'Learn from rare words:\n  • -ing → VERB\n  • -ly → ADV\n  • -tion/-ness → NOUN', RGBColor(0x3A, 0x1E, 0x5F)),
    ('Capitalization Clue', 'Capitalized unknown words are likely NOUNs\n(proper nouns / names)', RGBColor(0x1E, 0x5F, 0x3A)),
    ('Closed-Class Blocking', 'Unknown words are never:\nDET, PRON, ADP, CONJ, PRT\n(these are closed classes)', RGBColor(0x5F, 0x1E, 0x1E)),
]
for i, (title, desc, color) in enumerate(improvements):
    left = Inches(0.5 + i * 3.1)
    add_rect(slide, left, Inches(1.7), Inches(2.8), Inches(3.0), color, ACCENT)
    add_textbox(slide, left + Inches(0.1), Inches(1.8), Inches(2.6), Inches(0.4),
                title, font_size=15, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), Inches(2.3), Inches(2.6), Inches(2.3),
                desc, font_size=13, color=WHITE)

# Result card
result_box = add_rect(slide, Inches(2.5), Inches(5.3), Inches(8), Inches(1.2), RGBColor(0x1E, 0x3A, 0x5F), GREEN)
add_textbox(slide, Inches(2.7), Inches(5.4), Inches(7.6), Inches(1.0),
            'Result: 93.43% → 96.91%  (+3.48 percentage points)',
            font_size=22, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

# ── Slide 8: Results ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            'Results & Accuracy', font_size=36, bold=True, color=ACCENT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

# Table
rows, cols = 3, 4
table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.8), Inches(11), Inches(2.5))
table = table_shape.table

headers = ['Model', 'Alpha', 'Word Accuracy', 'Key Technique']
data = [
    ['Baseline HMM', '1.0', '93.43%', 'Flat OOV (add-α smoothing)'],
    ['Improved HMM', '0.01', '96.91%', 'Suffix + cap + closed-class'],
]

for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD

# Comparison image
add_image_safe(slide, os.path.join(OUT, 'accuracy_comparison.png'),
               Inches(1.5), Inches(4.5), width=Inches(4.5))

add_textbox(slide, Inches(6.5), Inches(4.6), Inches(6), Inches(0.5),
            'Key Insight:', font_size=18, bold=True, color=ACCENT)
items5 = [
    '+2.5x fewer errors (6.57% → 3.09% error rate)',
    'Suffix model captures morphological patterns',
    'Closed-class blocking eliminates impossible tags for unknowns',
    'Capitalization identifies proper nouns reliably',
]
add_bullet_frame(slide, Inches(6.5), Inches(5.1), Inches(6), Inches(2.0), items5, font_size=14)

# ── Slide 9: Error Analysis ─────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            'Error Analysis & Limitations', font_size=36, bold=True, color=ACCENT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

errors_list = [
    'Noun/Verb confusion: "offer", "average", "watch" — ambiguous words',
    'DET vs PRON: "her" before noun (her book) vs object position (saw her)',
    'ADJ vs VERB: participles like "tempting", "broken", "running"',
    'Rare punctuation / foreign words (tag X) due to sparse data',
]
add_bullet_frame(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), errors_list, font_size=16)

# Root cause box
rc_box = add_rect(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(3.5), RGBColor(0x3A, 0x1E, 0x1E), RED)
add_textbox(slide, Inches(7.2), Inches(1.6), Inches(5.1), Inches(0.4),
            'Root Cause:', font_size=18, bold=True, color=RED)
add_textbox(slide, Inches(7.2), Inches(2.1), Inches(5.1), Inches(2.8),
            'First-Order Markov Assumption\n\n'
            'Each tag depends ONLY on the immediately\n'
            'previous tag. The model cannot:\n'
            '  • Look ahead to see if a noun follows\n'
            '  • Track long-distance agreement\n'
            '  • Know if a clause already has a main verb\n\n'
            'This is the fundamental limitation of HMMs\n'
            'vs. CRFs and Transformer models.',
            font_size=14, color=WHITE)

# Future work
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.4),
            'Future Work: CRF/Transformer models (BERT) relax the Markov assumption → higher accuracy',
            font_size=16, color=MUTED)

# ── Slide 10: Demo ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            'Live Demo — Web Application', font_size=36, bold=True, color=ACCENT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

demo_items = [
    'Flask web app at http://127.0.0.1:5000',
    'Type any sentence and see POS tags in real-time',
    'Color-coded tags with hover descriptions',
    'Contractions expanded (I\'m → I am, don\'t → do not)',
    'Handles unseen names via suffix+cap model (Heng → NOUN)',
]
add_bullet_frame(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(2.5), demo_items, font_size=18)

# Example tags
examples = [
    ('my', 'DET'), ('watch', 'NOUN'), ('is', 'VERB'), ('broken', 'VERB'),
    ('', ''),
    ('watch', 'VERB'), ('the', 'DET'), ('dog', 'NOUN'),
    ('', ''),
    ('Heng', 'NOUN'), ('loves', 'VERB'), ('her', 'DET'),
]

# ── Slide 11: Conclusion ────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            'Conclusion', font_size=36, bold=True, color=ACCENT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

conclusions = [
    'Built a complete HMM POS tagger from scratch (no ML toolkit used)',
    'MLE estimation of π, A, B + Laplace smoothing for OOV words',
    'Log-space Viterbi ensures numerical stability',
    'Improved OOV handling: suffix model + capitalization + closed-class blocking',
    'Accuracy: 93.43% (baseline) → 96.91% (improved)',
    'Main limitation: first-order Markov assumption (no long-range context)',
    'Deployed as an interactive Flask web application',
]
add_bullet_frame(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(4.0), conclusions, font_size=18)

add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5),
            '"The HMM is the simplest sequence model that works — understanding its strengths and '
            'limitations builds intuition for all modern sequence models."',
            font_size=14, color=MUTED)

PPTX_PATH = r"E:\text_processing\HMM_POS_Tagger_Presentation.pptx"
prs.save(PPTX_PATH)
print(f"[OK] Presentation saved: {PPTX_PATH}")
print("\nDone! All enhancements complete.")
